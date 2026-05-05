"""
Falantir submission file generator.

Reads Falantir_Submission_Package.md and produces:
  - Falantir_Project_Report.docx  (Step B content)
  - Falantir_Project_Slides.pptx  (Step C content)

Usage:
  pip install python-docx python-pptx
  python generate_submission_files.py
"""

import re
from pathlib import Path

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt as PPTPt
from pptx.dml.color import RGBColor as PPTRGB

ROOT = Path(__file__).parent
SOURCE_MD = ROOT / "Falantir_Submission_Package.md"
DOCX_OUT = ROOT / "Falantir_Project_Report.docx"
PPTX_OUT = ROOT / "Falantir_Project_Slides.pptx"


def _split_sections(md: str):
    """Return (report_md, slide_outline_md) by splitting on the Step markers."""
    step_b_idx = md.index("# STEP B")
    step_c_idx = md.index("# STEP C")
    step_d_idx = md.index("# STEP D")
    report = md[step_b_idx:step_c_idx]
    slides = md[step_c_idx:step_d_idx]
    return report, slides


def _add_styled_heading(doc, text, level):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = RGBColor(0x1F, 0x29, 0x37)
    return h


def _parse_table(lines, start_idx):
    """Parse a Markdown pipe table starting at start_idx. Return (rows, end_idx)."""
    rows = []
    i = start_idx
    while i < len(lines) and lines[i].strip().startswith("|"):
        line = lines[i].strip().strip("|")
        cells = [c.strip() for c in line.split("|")]
        rows.append(cells)
        i += 1
    if len(rows) >= 2 and all(set(c) <= set("- :") for c in rows[1]):
        rows.pop(1)
    return rows, i


def build_docx(report_md: str):
    doc = Document()

    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    section = doc.sections[0]
    section.left_margin = Inches(1.0)
    section.right_margin = Inches(1.0)
    section.top_margin = Inches(1.0)
    section.bottom_margin = Inches(1.0)

    lines = report_md.splitlines()
    i = 0

    in_code = False
    code_buffer = []

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("```"):
            if in_code:
                p = doc.add_paragraph()
                run = p.add_run("\n".join(code_buffer))
                run.font.name = "Consolas"
                run.font.size = Pt(9)
                code_buffer = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_buffer.append(line)
            i += 1
            continue

        if not stripped:
            i += 1
            continue

        if stripped.startswith("# "):
            _add_styled_heading(doc, stripped[2:], level=0)
        elif stripped.startswith("## "):
            _add_styled_heading(doc, stripped[3:], level=1)
        elif stripped.startswith("### "):
            _add_styled_heading(doc, stripped[4:], level=2)
        elif stripped.startswith("#### "):
            _add_styled_heading(doc, stripped[5:], level=3)
        elif stripped == "---":
            doc.add_paragraph("_" * 60).alignment = WD_ALIGN_PARAGRAPH.CENTER
        elif stripped.startswith("|"):
            rows, new_i = _parse_table(lines, i)
            if rows:
                table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                table.style = "Light Grid Accent 1"
                for r_idx, row in enumerate(rows):
                    for c_idx, cell_text in enumerate(row):
                        if c_idx < len(table.rows[r_idx].cells):
                            cell = table.rows[r_idx].cells[c_idx]
                            cell.text = cell_text.replace("**", "")
                            for paragraph in cell.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(10)
            i = new_i
            continue
        elif stripped.startswith("- "):
            doc.add_paragraph(stripped[2:].replace("**", ""), style="List Bullet")
        elif re.match(r"^\d+\.\s", stripped):
            text = re.sub(r"^\d+\.\s", "", stripped)
            doc.add_paragraph(text.replace("**", ""), style="List Number")
        else:
            text = stripped
            text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
            text = re.sub(r"\*(.*?)\*", r"\1", text)
            text = re.sub(r"`(.*?)`", r"\1", text)
            doc.add_paragraph(text)
        i += 1

    doc.save(DOCX_OUT)
    print(f"Wrote {DOCX_OUT}")


def _parse_slides(slides_md: str):
    """Return list of dicts: [{title, key_points: [...], visual: str}, ...]"""
    slides = []
    blocks = re.split(r"^### Slide \d+ — ", slides_md, flags=re.MULTILINE)
    blocks = [b for b in blocks if b.strip()]
    for block in blocks[1:]:
        lines = block.splitlines()
        if not lines:
            continue
        slide_title = lines[0].strip()
        title = slide_title
        title_match = re.search(r"\*\*Title\*\*:\s*(.+)", block)
        if title_match:
            title = title_match.group(1).strip()
        key_points = []
        kp_match = re.search(r"\*\*Key Points\*\*:\s*\n((?:[-*].+\n?)+)", block)
        if kp_match:
            for line in kp_match.group(1).splitlines():
                line = line.strip()
                if line.startswith(("- ", "* ")):
                    key_points.append(line[2:].strip())
        visual = ""
        v_match = re.search(r"\*\*Suggested Visual\*\*:\s*(.+?)(?:\n\n|\n---|$)", block, re.DOTALL)
        if v_match:
            visual = v_match.group(1).strip()
        slides.append({
            "section_title": slide_title,
            "title": title,
            "key_points": key_points,
            "visual": visual,
        })
    return slides


def build_pptx(slides_md: str):
    prs = Presentation()
    prs.slide_width = PPTInches(13.333)
    prs.slide_height = PPTInches(7.5)

    slides_data = _parse_slides(slides_md)

    accent = PPTRGB(0x10, 0xB9, 0x81)
    dark = PPTRGB(0x0F, 0x17, 0x2A)
    light = PPTRGB(0xF8, 0xFA, 0xFC)

    for idx, slide_data in enumerate(slides_data, start=1):
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = light

        accent_bar = slide.shapes.add_shape(
            1, PPTInches(0), PPTInches(0), PPTInches(13.333), PPTInches(0.4)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent
        accent_bar.line.fill.background()
        accent_bar.text_frame.text = ""

        title_box = slide.shapes.add_textbox(
            PPTInches(0.5), PPTInches(0.6), PPTInches(12.3), PPTInches(1.0)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = slide_data["title"]
        run.font.size = PPTPt(32)
        run.font.bold = True
        run.font.color.rgb = dark

        if slide_data["key_points"]:
            content_box = slide.shapes.add_textbox(
                PPTInches(0.7), PPTInches(2.0), PPTInches(8.0), PPTInches(4.5)
            )
            ctf = content_box.text_frame
            ctf.word_wrap = True
            for j, point in enumerate(slide_data["key_points"]):
                if j == 0:
                    p = ctf.paragraphs[0]
                else:
                    p = ctf.add_paragraph()
                p.text = "•  " + point
                for run in p.runs:
                    run.font.size = PPTPt(18)
                    run.font.color.rgb = dark
                p.space_after = PPTPt(8)

        if slide_data["visual"]:
            visual_box = slide.shapes.add_textbox(
                PPTInches(9.0), PPTInches(2.0), PPTInches(4.0), PPTInches(4.5)
            )
            vtf = visual_box.text_frame
            vtf.word_wrap = True
            p = vtf.paragraphs[0]
            run = p.add_run()
            run.text = "Suggested Visual"
            run.font.size = PPTPt(14)
            run.font.bold = True
            run.font.color.rgb = accent
            p2 = vtf.add_paragraph()
            p2.text = slide_data["visual"]
            for r in p2.runs:
                r.font.size = PPTPt(11)
                r.font.italic = True
                r.font.color.rgb = dark

        page_box = slide.shapes.add_textbox(
            PPTInches(12.5), PPTInches(7.0), PPTInches(0.7), PPTInches(0.4)
        )
        ptf = page_box.text_frame
        pp = ptf.paragraphs[0]
        prun = pp.add_run()
        prun.text = f"{idx} / {len(slides_data)}"
        prun.font.size = PPTPt(10)
        prun.font.color.rgb = dark

    prs.save(PPTX_OUT)
    print(f"Wrote {PPTX_OUT}")


def main():
    if not SOURCE_MD.exists():
        raise SystemExit(f"Source not found: {SOURCE_MD}")
    md = SOURCE_MD.read_text(encoding="utf-8")
    report_md, slides_md = _split_sections(md)
    build_docx(report_md)
    build_pptx(slides_md)
    print("Done.")


if __name__ == "__main__":
    main()
